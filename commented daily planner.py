# you need to install python-pptx via - pip install python-pptx
# and of course you need Python...but no need for PowerPoint
# Importing the various references we use in the code

from datetime import date, timedelta
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

#function to loop through each date in date range we pass start and end date to it
def daterange(start_date, end_date):
    #and loop through the range of the amount of days between start and end date
    for n in range(int((end_date - start_date).days)):
        #then return the start date + the number in each loop
        yield start_date + timedelta(n)

#function to draw the shape & text that make up the main template
def drawdailyplanner():
    #refer to the new shape as line1, then if we need to make changes, add to a group etc
    #its easy to do
    #add a rounded rectangle, as it has more options than a connector shape or line
    #(left, top, width, height) you can put anything here and set below too. Will explain
    #this one and then leave the rest as they are the same just different locations!
    line1=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    #set the width as a percentage of the slide so that it scales on different sizes
    #int() return the number as an integer because if we get a decimal it will throw an error
    line1.width = int((presentation.slide_width / 100) * 59)
    line1.top = int((presentation.slide_height / 90) * 13)
    line1.left = int((presentation.slide_width / 90) * 28.8)
    #Set the shape height to 1pt to make our rounded rectangle appear as a line
    line1.height = Pt(1)
    #add a solid colour fill for our shape
    line1.fill.solid()
    #color our shape black
    line1.fill.fore_color.rgb = RGBColor(0, 0, 0)
    #remove any outside line colour on our rectangle
    line1.line.fill.background()
    #remove the shadow effect
    line1.shadow.inherit = False

    line2=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line2.width = int((presentation.slide_width / 100) * 22.8)
    line2.top = int((presentation.slide_height / 90) * 15.8)
    line2.left = int((presentation.slide_width / 100) * 9)
    line2.height = Pt(1)
    line2.fill.solid()
    line2.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line2.line.fill.background()
    line2.shadow.inherit = False

    line3=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line3.width = int((presentation.slide_width / 100) * 82)
    line3.top = int((presentation.slide_height / 90) * 18.6)
    line3.left = int((presentation.slide_width / 2) - (line3.width / 2))
    line3.height = Pt(1)
    line3.fill.solid()
    line3.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line3.line.fill.background()
    line3.shadow.inherit = False

    line4=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line4.width = int((presentation.slide_width / 100) * 22.8)
    line4.top = int((presentation.slide_height / 90) * 21.4)
    line4.left = int((presentation.slide_width / 100) * 9)
    line4.height = Pt(1)
    line4.fill.solid()
    line4.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line4.line.fill.background()
    line4.shadow.inherit = False

    line5=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line5.width = int((presentation.slide_width / 100) * 82)
    line5.top = int((presentation.slide_height / 90) * 24.2)
    line5.left = int((presentation.slide_width / 2) - (line5.width / 2))
    line5.height = Pt(1)
    line5.fill.solid()
    line5.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line5.line.fill.background()
    line5.shadow.inherit = False

    line6=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line6.width = int((presentation.slide_width / 100) * 22.8)
    line6.top = int((presentation.slide_height / 90) * 27)
    line6.left = int((presentation.slide_width / 100) * 9)
    line6.height = Pt(1)
    line6.fill.solid()
    line6.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line6.line.fill.background()
    line6.shadow.inherit = False

    line7=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line7.width = int((presentation.slide_width / 100) * 82)
    line7.top = int((presentation.slide_height / 90) * 29.8)
    line7.left = int((presentation.slide_width / 2) - (line7.width / 2))
    line7.height = Pt(1)
    line7.fill.solid()
    line7.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line7.line.fill.background()
    line7.shadow.inherit = False

    line8=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line8.width = int((presentation.slide_width / 100) * 22.8)
    line8.top = int((presentation.slide_height / 90) * 32.6)
    line8.left = int((presentation.slide_width / 100) * 9)
    line8.height = Pt(1)
    line8.fill.solid()
    line8.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line8.line.fill.background()
    line8.shadow.inherit = False

    line9=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line9.width = int((presentation.slide_width / 100) * 82)
    line9.top = int((presentation.slide_height / 90) * 35.4)
    line9.left = int((presentation.slide_width / 2) - (line9.width / 2))
    line9.height = Pt(1)
    line9.fill.solid()
    line9.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line9.line.fill.background()
    line9.shadow.inherit = False

    line10=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line10.width = int((presentation.slide_width / 100) * 22.8)
    line10.top = int((presentation.slide_height / 90) * 38.2)
    line10.left = int((presentation.slide_width / 100) * 9)
    line10.height = Pt(1)
    line10.fill.solid()
    line10.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line10.line.fill.background()
    line10.shadow.inherit = False

    line11=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line11.width = int((presentation.slide_width / 100) * 82)
    line11.top = int((presentation.slide_height / 90) * 41)
    line11.left = int((presentation.slide_width / 2) - (line11.width / 2))
    line11.height = Pt(1)
    line11.fill.solid()
    line11.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line11.line.fill.background()
    line11.shadow.inherit = False

    line12=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line12.width = int((presentation.slide_width / 100) * 22.8)
    line12.top = int((presentation.slide_height / 90) * 43.8)
    line12.left = int((presentation.slide_width / 100) * 9)
    line12.height = Pt(1)
    line12.fill.solid()
    line12.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line12.line.fill.background()
    line12.shadow.inherit = False

    line13=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line13.width = int((presentation.slide_width / 100) * 82)
    line13.top = int((presentation.slide_height / 90) * 46.6)
    line13.left = int((presentation.slide_width / 2) - (line13.width / 2))
    line13.height = Pt(1)
    line13.fill.solid()
    line13.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line13.line.fill.background()
    line13.shadow.inherit = False

    line14=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line14.width = int((presentation.slide_width / 100) * 59)
    line14.top = int((presentation.slide_height / 90) * 52.2)
    line14.left = int((presentation.slide_width / 90) * 28.8)
    line14.height = Pt(1)
    line14.fill.solid()
    line14.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line14.line.fill.background()
    line14.shadow.inherit = False

    line15=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line15.width = int((presentation.slide_width / 100) * 22.8)
    line15.top = int((presentation.slide_height / 90) * 55)
    line15.left = int((presentation.slide_width / 100) * 9)
    line15.height = Pt(1)
    line15.fill.solid()
    line15.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line15.line.fill.background()
    line15.shadow.inherit = False

    line16=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line16.width = int((presentation.slide_width / 100) * 82)
    line16.top = int((presentation.slide_height / 90) * 57.8)
    line16.left = int((presentation.slide_width / 2) - (line16.width / 2))
    line16.height = Pt(1)
    line16.fill.solid()
    line16.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line16.line.fill.background()
    line16.shadow.inherit = False

    line17=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line17.width = int((presentation.slide_width / 100) * 22.8)
    line17.top = int((presentation.slide_height / 90) * 60.6)
    line17.left = int((presentation.slide_width / 100) * 9)
    line17.height = Pt(1)
    line17.fill.solid()
    line17.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line17.line.fill.background()
    line17.shadow.inherit = False

    line18=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line18.width = int((presentation.slide_width / 100) * 82)
    line18.top = int((presentation.slide_height / 90) * 63.4)
    line18.left = int((presentation.slide_width / 2) - (line18.width / 2))
    line18.height = Pt(1)
    line18.fill.solid()
    line18.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line18.line.fill.background()
    line18.shadow.inherit = False

    line19=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line19.width = int((presentation.slide_width / 100) * 22.8)
    line19.top = int((presentation.slide_height / 90) * 66.2)
    line19.left = int((presentation.slide_width / 100) * 9)
    line19.height = Pt(1)
    line19.fill.solid()
    line19.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line19.line.fill.background()
    line19.shadow.inherit = False

    line20=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line20.width = int((presentation.slide_width / 100) * 82)
    line20.top = int((presentation.slide_height / 90) * 69)
    line20.left = int((presentation.slide_width / 2) - (line20.width / 2))
    line20.height = Pt(1)
    line20.fill.solid()
    line20.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line20.line.fill.background()
    line20.shadow.inherit = False

    line21=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line21.width = int((presentation.slide_width / 100) * 22.8)
    line21.top = int((presentation.slide_height / 90) * 71.8)
    line21.left = int((presentation.slide_width / 100) * 9)
    line21.height = Pt(1)
    line21.fill.solid()
    line21.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line21.line.fill.background()
    line21.shadow.inherit = False

    line22=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line22.width = int((presentation.slide_width / 100) * 82)
    line22.top = int((presentation.slide_height / 90) * 74.6)
    line22.left = int((presentation.slide_width / 2) - (line22.width / 2))
    line22.height = Pt(1)
    line22.fill.solid()
    line22.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line22.line.fill.background()
    line22.shadow.inherit = False

    line23=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line23.width = int((presentation.slide_width / 100) * 22.8)
    line23.top = int((presentation.slide_height / 90) * 77.4)
    line23.left = int((presentation.slide_width / 100) * 9)
    line23.height = Pt(1)
    line23.fill.solid()
    line23.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line23.line.fill.background()
    line23.shadow.inherit = False

    line24=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line24.width = int((presentation.slide_width / 100) * 82)
    line24.top = int((presentation.slide_height / 90) * 80.2)
    line24.left = int((presentation.slide_width / 2) - (line24.width / 2))
    line24.height = Pt(1)
    line24.fill.solid()
    line24.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line24.line.fill.background()
    line24.shadow.inherit = False

    line25=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line25.width = int((presentation.slide_width / 100) * 22.8)
    line25.top = int((presentation.slide_height / 90) * 83)
    line25.left = int((presentation.slide_width / 100) * 9)
    line25.height = Pt(1)
    line25.fill.solid()
    line25.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line25.line.fill.background()
    line25.shadow.inherit = False

    line26=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line26.width = int((presentation.slide_width / 100) * 82)
    line26.top = int((presentation.slide_height / 90) * 85.8)
    line26.left = int((presentation.slide_width / 2) - (line22.width / 2))
    line26.height = Pt(1)
    line26.fill.solid()
    line26.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line26.line.fill.background()
    line26.shadow.inherit = False

    # other direction

    line27=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line27.height = int((presentation.slide_height / 100) * 80.8)
    line27.top = int((presentation.slide_height / 90) * 13)
    line27.left = int((presentation.slide_width / 90) * 28.8)
    line27.width = Pt(1)
    line27.fill.solid()
    line27.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line27.line.fill.background()
    line27.shadow.inherit = False

    line28=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line28.height = int((presentation.slide_height / 100) * 74.6)
    line28.top = int((presentation.slide_height / 90) * 13)
    line28.left = int((presentation.slide_width / 90) * 33.6)
    line28.width = Pt(1)
    line28.fill.solid()
    line28.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line28.line.fill.background()
    line28.shadow.inherit = False

    line29=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line29.height = int((presentation.slide_height / 100) * 80.8)
    line29.top = int((presentation.slide_height / 90) * 13)
    line29.left = int((presentation.slide_width / 100) * 91)
    line29.width = Pt(1)
    line29.fill.solid()
    line29.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line29.line.fill.background()
    line29.shadow.inherit = False

    line30=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line30.height = int((presentation.slide_height / 100) * 74.6)
    line30.top = int((presentation.slide_height / 90) * 13)
    line30.left = int((presentation.slide_width / 90) * 55.35)
    line30.width = Pt(1)
    line30.fill.solid()
    line30.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line30.line.fill.background()
    line30.shadow.inherit = False

    line31=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(2), Inches(1), Inches(2))
    line31.height = int((presentation.slide_height / 100) * 74.6)
    line31.top = int((presentation.slide_height / 90) * 13)
    line31.left = int((presentation.slide_width / 90) * 60.15)
    line31.width = Pt(1)
    line31.fill.solid()
    line31.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line31.line.fill.background()
    line31.shadow.inherit = False
    
    #refer to the new textbox as txBox2, then if we need to make changes, add to a group etc
    #its easy to do
    #add a textbox
    #(left, top, width, height) you can put anything here and set below too. Will explain
    #this one and then leave the rest as they are the same just different locations!
    txBox2 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    #set the left as a percentage of the slide so that it scales on different sizes
    #int() return the number as an integer because if we get a decimal it will throw an error
    txBox2.left = int((presentation.slide_width / 100) * 9)
    txBox2.width = int((presentation.slide_width / 100) * 22.8)
    txBox2.height = Inches(0.27777777777)
    #using shorthand to refer to elements to save long winded references
    tf2 = txBox2.text_frame
    #using shorthand to refer to elements to save long winded references
    p2 = tf2.paragraphs[0]
    #set our text to Priority Tasks
    p2.text = "Priority Tasks"
    #Set our font size as 20pt
    p2.font.size = Pt(20)
    #Align our text in the centre
    p2.alignment = PP_ALIGN.CENTER
    txBox2.top = int(((presentation.slide_height / 90) * 13) - txBox2.height)

    txBox3 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox3.left = int((presentation.slide_width / 100) * 9)
    txBox3.width = int((presentation.slide_width / 100) * 22.8)
    txBox3.height = Inches(0.27777777777)
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "To Do"
    p3.font.size = Pt(20)
    p3.alignment = PP_ALIGN.CENTER
    txBox3.top = int(((presentation.slide_height / 90) * 52.2) - txBox3.height)

    txBox4 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox4.left = int((presentation.slide_width / 90) * 33.6)
    txBox4.width = int((presentation.slide_width / 90) * 21.75)
    txBox4.height = Inches(0.27777777777)
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "AM"
    p4.font.size = Pt(20)
    p4.alignment = PP_ALIGN.CENTER
    txBox4.top = int(((presentation.slide_height / 90) * 12) - txBox4.height)

    txBox5 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox5.left = int((presentation.slide_width / 90) * 60.15)
    txBox5.width = int((presentation.slide_width / 90) * 21.75)
    txBox5.height = Inches(0.27777777777)
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "PM"
    p5.font.size = Pt(20)
    p5.alignment = PP_ALIGN.CENTER
    txBox5.top = int(((presentation.slide_height / 90) * 12) - txBox4.height)

    txBox6 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox6.left = int((presentation.slide_width / 100) * 32)
    txBox6.width = int((presentation.slide_width / 100) * 59)
    txBox6.height = int((presentation.slide_height / 90) * 5.6)
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "WATER INTAKE - 1 2 3 4 5 6 7 8"
    p6.font.size = Pt(20)
    p6.alignment = PP_ALIGN.CENTER
    txBox6.top = int(((presentation.slide_height / 90) * 80.2))

    # left column  times

    txBox7 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox7.left = int((presentation.slide_width / 90) * 28.8)
    txBox7.width = int((presentation.slide_width / 90) * 4.8)
    txBox7.height = int((presentation.slide_height / 90) * 5.6)
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "12"
    p7.font.size = Pt(16)
    p7.alignment = PP_ALIGN.CENTER
    txBox7.top = int((presentation.slide_height / 90) * 13)

    txBox8 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox8.left = int((presentation.slide_width / 90) * 28.8)
    txBox8.width = int((presentation.slide_width / 90) * 4.8)
    txBox8.height = int((presentation.slide_height / 90) * 5.6)
    tf8 = txBox8.text_frame
    p8 = tf8.paragraphs[0]
    p8.text = "1"
    p8.font.size = Pt(16)
    p8.alignment = PP_ALIGN.CENTER
    txBox8.top = int((presentation.slide_height / 90) * 18.6)

    txBox9 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox9.left = int((presentation.slide_width / 90) * 28.8)
    txBox9.width = int((presentation.slide_width / 90) * 4.8)
    txBox9.height = int((presentation.slide_height / 90) * 5.6)
    tf9 = txBox9.text_frame
    p9 = tf9.paragraphs[0]
    p9.text = "2"
    p9.font.size = Pt(16)
    p9.alignment = PP_ALIGN.CENTER
    txBox9.top = int((presentation.slide_height / 90) * 24.2)

    txBox10 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox10.left = int((presentation.slide_width / 90) * 28.8)
    txBox10.width = int((presentation.slide_width / 90) * 4.8)
    txBox10.height = int((presentation.slide_height / 90) * 5.6)
    tf10 = txBox10.text_frame
    p10 = tf10.paragraphs[0]
    p10.text = "3"
    p10.font.size = Pt(16)
    p10.alignment = PP_ALIGN.CENTER
    txBox10.top = int((presentation.slide_height / 90) * 29.8)

    txBox11 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox11.left = int((presentation.slide_width / 90) * 28.8)
    txBox11.width = int((presentation.slide_width / 90) * 4.8)
    txBox11.height = int((presentation.slide_height / 90) * 5.6)
    tf11 = txBox11.text_frame
    p11 = tf11.paragraphs[0]
    p11.text = "4"
    p11.font.size = Pt(16)
    p11.alignment = PP_ALIGN.CENTER
    txBox11.top = int((presentation.slide_height / 90) * 35.4)

    txBox12 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox12.left = int((presentation.slide_width / 90) * 28.8)
    txBox12.width = int((presentation.slide_width / 90) * 4.8)
    txBox12.height = int((presentation.slide_height / 90) * 5.6)
    tf12 = txBox12.text_frame
    p12 = tf12.paragraphs[0]
    p12.text = "5"
    p12.font.size = Pt(16)
    p12.alignment = PP_ALIGN.CENTER
    txBox12.top = int((presentation.slide_height / 90) * 41)

    txBox13 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox13.left = int((presentation.slide_width / 90) * 28.8)
    txBox13.width = int((presentation.slide_width / 90) * 4.8)
    txBox13.height = int((presentation.slide_height / 90) * 5.6)
    tf13 = txBox13.text_frame
    p13 = tf13.paragraphs[0]
    p13.text = "6"
    p13.font.size = Pt(16)
    p13.alignment = PP_ALIGN.CENTER
    txBox13.top = int((presentation.slide_height / 90) * 46.6)

    txBox14 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox14.left = int((presentation.slide_width / 90) * 28.8)
    txBox14.width = int((presentation.slide_width / 90) * 4.8)
    txBox14.height = int((presentation.slide_height / 90) * 5.6)
    tf14 = txBox14.text_frame
    p14 = tf14.paragraphs[0]
    p14.text = "7"
    p14.font.size = Pt(16)
    p14.alignment = PP_ALIGN.CENTER
    txBox14.top = int((presentation.slide_height / 90) * 52.2)

    txBox15 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox15.left = int((presentation.slide_width / 90) * 28.8)
    txBox15.width = int((presentation.slide_width / 90) * 4.8)
    txBox15.height = int((presentation.slide_height / 90) * 5.6)
    tf15 = txBox15.text_frame
    p15 = tf15.paragraphs[0]
    p15.text = "8"
    p15.font.size = Pt(16)
    p15.alignment = PP_ALIGN.CENTER
    txBox15.top = int((presentation.slide_height / 90) * 57.8)

    txBox16 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox16.left = int((presentation.slide_width / 90) * 28.8)
    txBox16.width = int((presentation.slide_width / 90) * 4.8)
    txBox16.height = int((presentation.slide_height / 90) * 5.6)
    tf16 = txBox16.text_frame
    p16 = tf16.paragraphs[0]
    p16.text = "9"
    p16.font.size = Pt(16)
    p16.alignment = PP_ALIGN.CENTER
    txBox16.top = int((presentation.slide_height / 90) * 63.4)

    txBox17 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox17.left = int((presentation.slide_width / 90) * 28.8)
    txBox17.width = int((presentation.slide_width / 90) * 4.8)
    txBox17.height = int((presentation.slide_height / 90) * 5.6)
    tf17 = txBox17.text_frame
    p17 = tf17.paragraphs[0]
    p17.text = "10"
    p17.font.size = Pt(16)
    p17.alignment = PP_ALIGN.CENTER
    txBox17.top = int((presentation.slide_height / 90) * 69)

    txBox18 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox18.left = int((presentation.slide_width / 90) * 28.8)
    txBox18.width = int((presentation.slide_width / 90) * 4.8)
    txBox18.height = int((presentation.slide_height / 90) * 5.6)
    tf18 = txBox18.text_frame
    p18 = tf18.paragraphs[0]
    p18.text = "11"
    p18.font.size = Pt(16)
    p18.alignment = PP_ALIGN.CENTER
    txBox18.top = int((presentation.slide_height / 90) * 74.6)

    #right column times

    txBox19 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox19.left = int((presentation.slide_width / 90) * 55.35)
    txBox19.width = int((presentation.slide_width / 90) * 4.8)
    txBox19.height = int((presentation.slide_height / 90) * 5.6)
    tf19 = txBox19.text_frame
    p19 = tf19.paragraphs[0]
    p19.text = "12"
    p19.font.size = Pt(16)
    p19.alignment = PP_ALIGN.CENTER
    txBox19.top = int((presentation.slide_height / 90) * 13)

    txBox20 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox20.left = int((presentation.slide_width / 90) * 55.35)
    txBox20.width = int((presentation.slide_width / 90) * 4.8)
    txBox20.height = int((presentation.slide_height / 90) * 5.6)
    tf20 = txBox20.text_frame
    p20 = tf20.paragraphs[0]
    p20.text = "1"
    p20.font.size = Pt(16)
    p20.alignment = PP_ALIGN.CENTER
    txBox20.top = int((presentation.slide_height / 90) * 18.6)

    txBox21 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox21.left = int((presentation.slide_width / 90) * 55.35)
    txBox21.width = int((presentation.slide_width / 90) * 4.8)
    txBox21.height = int((presentation.slide_height / 90) * 5.6)
    tf21 = txBox21.text_frame
    p21 = tf21.paragraphs[0]
    p21.text = "2"
    p21.font.size = Pt(16)
    p21.alignment = PP_ALIGN.CENTER
    txBox21.top = int((presentation.slide_height / 90) * 24.2)

    txBox22 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox22.left = int((presentation.slide_width / 90) * 55.35)
    txBox22.width = int((presentation.slide_width / 90) * 4.8)
    txBox22.height = int((presentation.slide_height / 90) * 5.6)
    tf22 = txBox22.text_frame
    p22 = tf22.paragraphs[0]
    p22.text = "3"
    p22.font.size = Pt(16)
    p22.alignment = PP_ALIGN.CENTER
    txBox22.top = int((presentation.slide_height / 90) * 29.8)

    txBox23 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox23.left = int((presentation.slide_width / 90) * 55.35)
    txBox23.width = int((presentation.slide_width / 90) * 4.8)
    txBox23.height = int((presentation.slide_height / 90) * 5.6)
    tf23 = txBox23.text_frame
    p23 = tf23.paragraphs[0]
    p23.text = "4"
    p23.font.size = Pt(16)
    p23.alignment = PP_ALIGN.CENTER
    txBox23.top = int((presentation.slide_height / 90) * 35.4)

    txBox24 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox24.left = int((presentation.slide_width / 90) * 55.35)
    txBox24.width = int((presentation.slide_width / 90) * 4.8)
    txBox24.height = int((presentation.slide_height / 90) * 5.6)
    tf24 = txBox24.text_frame
    p24 = tf24.paragraphs[0]
    p24.text = "5"
    p24.font.size = Pt(16)
    p24.alignment = PP_ALIGN.CENTER
    txBox24.top = int((presentation.slide_height / 90) * 41)

    txBox25 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox25.left = int((presentation.slide_width / 90) * 55.35)
    txBox25.width = int((presentation.slide_width / 90) * 4.8)
    txBox25.height = int((presentation.slide_height / 90) * 5.6)
    tf25 = txBox25.text_frame
    p25 = tf25.paragraphs[0]
    p25.text = "6"
    p25.font.size = Pt(16)
    p25.alignment = PP_ALIGN.CENTER
    txBox25.top = int((presentation.slide_height / 90) * 46.6)

    txBox26 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox26.left = int((presentation.slide_width / 90) * 55.35)
    txBox26.width = int((presentation.slide_width / 90) * 4.8)
    txBox26.height = int((presentation.slide_height / 90) * 5.6)
    tf26 = txBox26.text_frame
    p26 = tf26.paragraphs[0]
    p26.text = "7"
    p26.font.size = Pt(16)
    p26.alignment = PP_ALIGN.CENTER
    txBox26.top = int((presentation.slide_height / 90) * 52.2)

    txBox27 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox27.left = int((presentation.slide_width / 90) * 55.35)
    txBox27.width = int((presentation.slide_width / 90) * 4.8)
    txBox27.height = int((presentation.slide_height / 90) * 5.6)
    tf27 = txBox27.text_frame
    p27 = tf27.paragraphs[0]
    p27.text = "8"
    p27.font.size = Pt(16)
    p27.alignment = PP_ALIGN.CENTER
    txBox27.top = int((presentation.slide_height / 90) * 57.8)

    txBox28 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox28.left = int((presentation.slide_width / 90) * 55.35)
    txBox28.width = int((presentation.slide_width / 90) * 4.8)
    txBox28.height = int((presentation.slide_height / 90) * 5.6)
    tf28 = txBox28.text_frame
    p28 = tf28.paragraphs[0]
    p28.text = "9"
    p28.font.size = Pt(16)
    p28.alignment = PP_ALIGN.CENTER
    txBox28.top = int((presentation.slide_height / 90) * 63.4)

    txBox29 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox29.left = int((presentation.slide_width / 90) * 55.35)
    txBox29.width = int((presentation.slide_width / 90) * 4.8)
    txBox29.height = int((presentation.slide_height / 90) * 5.6)
    tf29 = txBox29.text_frame
    p29 = tf29.paragraphs[0]
    p29.text = "10"
    p29.font.size = Pt(16)
    p29.alignment = PP_ALIGN.CENTER
    txBox29.top = int((presentation.slide_height / 90) * 69)

    txBox30 = slide.shapes.add_textbox((0 + Inches(2)), (0 + Inches(2)), (presentation.slide_width), Inches(1))
    txBox30.left = int((presentation.slide_width / 90) * 55.35)
    txBox30.width = int((presentation.slide_width / 90) * 4.8)
    txBox30.height = int((presentation.slide_height / 90) * 5.6)
    tf30 = txBox30.text_frame
    p30 = tf30.paragraphs[0]
    p30.text = "11"
    p30.font.size = Pt(16)
    p30.alignment = PP_ALIGN.CENTER
    txBox30.top = int((presentation.slide_height / 90) * 74.6)


    



    
#create a new presentation to create our daily planner in
presentation = Presentation()

#set our presentation slide width & height
presentation.slide_width = Inches(11)
presentation.slide_height = Inches(11)

#blank layout without any textbox placeholders
layout = presentation.slide_masters[0].slide_layouts[6]

#start and end date for our daily planner
start_date = date(2022, 1, 1)
end_date = date(2023, 1, 1)

#loop through our date range and... 
for single_date in daterange(start_date, end_date):

    #add a new slide with the blank layout we setup earlier
    slide = presentation.slides.add_slide(layout)

    #refer to new text box as txBox which makes it easy to refer to in our code or change etc.
    #add a textbox (left, top, width, height) you can put anything here and set below too
    txBox = slide.shapes.add_textbox((0 + Inches(1)), (0 + Inches(0.1)), (presentation.slide_width), Inches(0.5))
    #changing the left location
    txBox.left = 0
    #changing the height
    txBox.height = Inches(0.27777777777)
    #changing the top location
    #when placing items if we want them to scale we need to be fluid so when we
    #change the slide size their location also adjusts so this is done as a fraction
    #of the overal slide size. As we may get a decimal result which would return an
    #error we use int() to make sure the result is a whole integer
    txBox.top = int((presentation.slide_height - 20) / 32)
    #using shorthand to refer to elements to save long winded references
    tf = txBox.text_frame
    #using shorthand to refer to elements to save long winded references
    p = tf.paragraphs[0]
    #the text for our textbox is the current date in our date loop with the format
    #%A %d %B %Y you can read all the formats you can choose from the link below
    # https://mkaz.blog/code/python-dates/
    p.text = single_date.strftime("%A %d %B %Y")
    #our font size
    p.font.size = Pt(28)
    #text alignment
    p.alignment = PP_ALIGN.CENTER

   


    #run our function (at the top)
    drawdailyplanner()
    
#save our completed presentation
presentation.save("C:\\Users\\lukeb\\Downloads\\test\\dates.pptx")
